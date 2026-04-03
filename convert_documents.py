#!/usr/bin/env python3
"""
Script de conversion des documents Édu-Connect
Markdown → PDF, Word, PowerPoint
"""

import os
from fpdf import FPDF
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN

# Fichiers à convertir
files_to_convert = [
    ('EDU_CONNECT_PRESENTATION_OFFICIELLE.md', 'Présentation Complète'),
    ('EDU_CONNECT_RESUME_EXECUTIF.md', 'Résumé Exécutif'),
    ('GUIDE_PRESENTATION_POWERPOINT.md', 'Guide PowerPoint')
]

class PDF(FPDF):
    def header(self):
        # Logo (si disponible)
        self.set_font('Arial', 'B', 15)
        self.cell(0, 10, 'Édu-Connect - Ministère de l\'Éducation RDC', 0, 1, 'C')
        self.ln(5)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

    def chapter_title(self, title):
        self.set_font('Arial', 'B', 14)
        self.cell(0, 10, title, 0, 1, 'L')
        self.ln(3)

    def chapter_body(self, body):
        self.set_font('Arial', '', 11)
        self.multi_cell(0, 6, body)
        self.ln()

def convert_to_pdf(md_file, output_file):
    """Convertir Markdown en PDF"""
    print(f"📄 Conversion PDF : {md_file}...")
    
    pdf = PDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Fonction pour nettoyer le texte
    def clean_text(text):
        # Remplacer les emojis par du texte
        replacements = {
            '📊': '[STATS]', '🎯': '[OBJECTIF]', '✅': '[OK]', '❌': '[NON]',
            '🔴': '[P0]', '🟡': '[P1]', '🟢': '[P2]', '🔵': '[INFO]',
            '⚠️': '[ATTENTION]', '💰': '[FINANCE]', '📋': '[PLAN]',
            '👀': '[OBSERVATION]', '📄': '[DOCUMENT]', '🚧': '[EN COURS]',
            '🏆': '[SUCCES]', '💻': '[TECH]', '🗺️': '[CARTE]', '📅': '[DATE]',
            '📝': '[NOTE]', '🎓': '[EDUCATION]', '👨‍🏫': '[ENSEIGNANT]',
            '👨‍🎓': '[ELEVE]', '🏫': '[ECOLE]', '🇨🇩': '[RDC]',
            '🚀': '[ACTION]', '📞': '[CONTACT]', '💡': '[IDEE]',
            '📈': '[CROISSANCE]', '⚡': '[PERFORMANCE]', '🔐': '[SECURITE]',
            '📱': '[MOBILE]', '🤖': '[IA]', '💬': '[MESSAGE]', '🎨': '[DESIGN]',
            '🌐': '[WEB]', '🔗': '[LIEN]', '📚': '[BIBLIOTHEQUE]', '⭐': '[STAR]'
        }
        for emoji, replacement in replacements.items():
            text = text.replace(emoji, replacement)
        
        # Nettoyer le markdown
        text = text.replace('**', '').replace('*', '').replace('`', '')
        text = text.replace('###', '').replace('##', '').replace('#', '')
        
        # Garder uniquement les caractères ASCII étendus
        text = ''.join(c if ord(c) < 256 else ' ' for c in text)
        return text.strip()
    
    # Traiter le contenu ligne par ligne
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        
        if not line:
            continue
            
        try:
            # Titres
            if line.startswith('# '):
                pdf.set_font('Arial', 'B', 16)
                clean_line = clean_text(line.replace('# ', ''))
                if clean_line:
                    pdf.cell(0, 10, clean_line[:150], 0, 1, 'L')
                    pdf.ln(2)
            elif line.startswith('## '):
                pdf.set_font('Arial', 'B', 14)
                clean_line = clean_text(line.replace('## ', ''))
                if clean_line:
                    pdf.cell(0, 8, clean_line[:150], 0, 1, 'L')
                    pdf.ln(2)
            elif line.startswith('### '):
                pdf.set_font('Arial', 'B', 12)
                clean_line = clean_text(line.replace('### ', ''))
                if clean_line:
                    pdf.cell(0, 7, clean_line[:150], 0, 1, 'L')
                    pdf.ln(1)
            # Texte normal
            else:
                clean_line = clean_text(line)
                if clean_line and len(clean_line) > 1:
                    pdf.set_font('Arial', '', 10)
                    pdf.multi_cell(0, 5, clean_line[:500])
        except Exception as e:
            # Ignorer les lignes problématiques
            continue
    
    pdf.output(output_file)
    print(f"✅ PDF créé : {output_file}")

def convert_to_word(md_file, output_file):
    """Convertir Markdown en Word"""
    print(f"📝 Conversion Word : {md_file}...")
    
    doc = Document()
    
    # Style du document
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(11)
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        
        if not line:
            continue
        
        # Titres
        if line.startswith('# '):
            p = doc.add_heading(line.replace('# ', ''), level=1)
        elif line.startswith('## '):
            p = doc.add_heading(line.replace('## ', ''), level=2)
        elif line.startswith('### '):
            p = doc.add_heading(line.replace('### ', ''), level=3)
        elif line.startswith('#### '):
            p = doc.add_heading(line.replace('#### ', ''), level=4)
        # Listes
        elif line.startswith('- ') or line.startswith('* '):
            clean_line = line[2:].replace('**', '').replace('`', '')
            doc.add_paragraph(clean_line, style='List Bullet')
        # Texte normal
        else:
            clean_line = line.replace('**', '').replace('*', '').replace('`', '')
            if clean_line:
                doc.add_paragraph(clean_line)
    
    doc.save(output_file)
    print(f"✅ Word créé : {output_file}")

def create_powerpoint():
    """Créer PowerPoint de présentation"""
    print(f"📊 Création PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 : Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Édu-Connect"
    subtitle.text = "Plateforme Éducative Nationale\nRépublique Démocratique du Congo\n\nMinistère de l'Éducation nationale et de la Nouvelle Citoyenneté"
    
    # Slide 2 : Vision & Mission
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Vision & Mission"
    body = slide.placeholders[1].text_frame
    body.text = "Vision : Digitaliser l'éducation en RDC\n"
    body.text += "Mission : Centraliser, Digitaliser, Optimiser\n"
    body.text += "Objectif : Conformité internationale (SIGE)"
    
    # Slide 3 : Chiffres Clés
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Chiffres Clés"
    body = slide.placeholders[1].text_frame
    body.text = "300 Établissements\n"
    body.text += "10 Enseignants\n"
    body.text += "20 Élèves\n"
    body.text += "600 Classes\n"
    body.text += "26 Provinces Administratives\n"
    body.text += "16 Profils Utilisateurs\n"
    body.text += "90-95% de Complétion"
    
    # Slide 4 : Modules Fonctionnels
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "20 Modules Fonctionnels"
    body = slide.placeholders[1].text_frame
    body.text = "Gestion Administrative\n"
    body.text += "Gestion Pédagogique\n"
    body.text += "Gestion Enseignants\n"
    body.text += "SIRH/DINACOPE\n"
    body.text += "GED (Documents)\n"
    body.text += "Statistiques & Rapports\n"
    body.text += "APIs Externes (Modules 3 & 4)\n"
    body.text += "Carte Scolaire Numérique"
    
    # Slide 5 : Conformité SIGE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conformité SIGE : 90-95%"
    body = slide.placeholders[1].text_frame
    body.text = "SIGE Unique : 100%\n"
    body.text += "SIGE Décentralisé : 100%\n"
    body.text += "SIGE Basé TIC : 100%\n"
    body.text += "SIGE Pérenne : 75%\n\n"
    body.text += "SCORE GLOBAL : 90-95%"
    
    # Slide 6 : Modules à Développer
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "4 Modules à Développer"
    body = slide.placeholders[1].text_frame
    body.text = "P0 - Module Budgétaire (3-4 sem, $3-4K)\n"
    body.text += "P0 - Plan d'Opérations (2-3 sem, $2-3K)\n"
    body.text += "P0 - Observations Leçons (2-3 sem, $2-3K)\n"
    body.text += "P1 - Exports Avancés (2 sem, $1.5-2K)"
    
    # Slide 7 : Budget
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Budget Année 1"
    body = slide.placeholders[1].text_frame
    body.text = "Hébergement Cloud : $600\n"
    body.text += "Module Budgétaire : $3,000-4,000\n"
    body.text += "Plan d'Opérations : $2,000-3,000\n"
    body.text += "Observations Leçons : $2,000-3,000\n"
    body.text += "Exports Avancés : $1,500-2,000\n"
    body.text += "Carte GPS : $500-1,000\n\n"
    body.text += "TOTAL : $9,900-13,600"
    
    # Slide 8 : Planning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Planning de Déploiement"
    body = slide.placeholders[1].text_frame
    body.text = "Phase 1 : Déploiement (Sem. 1-2)\n"
    body.text += "Phase 2 : Modules Critiques (Mois 1-3)\n"
    body.text += "Phase 3 : Formation Nationale (Mois 3-6)\n"
    body.text += "Phase 4 : Amélioration Continue (Mois 6-12)"
    
    # Slide 9 : Recommandation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Recommandation Stratégique"
    body = slide.placeholders[1].text_frame
    body.text = "Déploiement Immédiat\n"
    body.text += "+ Développement Modules Critiques\n\n"
    body.text += "Justification :\n"
    body.text += "- Système mature (90-95%)\n"
    body.text += "- Conforme SIGE (90-95%)\n"
    body.text += "- Technologie pérenne (10+ ans)\n"
    body.text += "- Investissement raisonnable (~$10-14K)"
    
    # Slide 10 : Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    body = slide.placeholders[1].text_frame
    body.text = "Édu-Connect est PRÊTE\n\n"
    body.text += "20 modules opérationnels\n"
    body.text += "90-95% conforme SIGE\n"
    body.text += "Technologie moderne et pérenne\n"
    body.text += "100% adapté à la RDC\n\n"
    body.text += "Pour une Éducation Digitale\n"
    body.text += "et Performante en RDC"
    
    output_file = '/app/EDU_CONNECT_PRESENTATION.pptx'
    prs.save(output_file)
    print(f"✅ PowerPoint créé : {output_file}")

# Exécution
if __name__ == '__main__':
    print("🚀 Début de la conversion des documents Édu-Connect...\n")
    
    # Convertir en PDF
    for md_file, name in files_to_convert:
        md_path = f'/app/{md_file}'
        pdf_path = f'/app/{md_file.replace(".md", ".pdf")}'
        
        if os.path.exists(md_path):
            convert_to_pdf(md_path, pdf_path)
        else:
            print(f"⚠️ Fichier non trouvé : {md_path}")
    
    print("\n")
    
    # Convertir en Word
    for md_file, name in files_to_convert:
        md_path = f'/app/{md_file}'
        docx_path = f'/app/{md_file.replace(".md", ".docx")}'
        
        if os.path.exists(md_path):
            convert_to_word(md_path, docx_path)
        else:
            print(f"⚠️ Fichier non trouvé : {md_path}")
    
    print("\n")
    
    # Créer PowerPoint
    create_powerpoint()
    
    print("\n✅ CONVERSION TERMINÉE !")
    print("\n📁 Fichiers créés :")
    print("   - EDU_CONNECT_PRESENTATION_OFFICIELLE.pdf")
    print("   - EDU_CONNECT_PRESENTATION_OFFICIELLE.docx")
    print("   - EDU_CONNECT_RESUME_EXECUTIF.pdf")
    print("   - EDU_CONNECT_RESUME_EXECUTIF.docx")
    print("   - GUIDE_PRESENTATION_POWERPOINT.pdf")
    print("   - GUIDE_PRESENTATION_POWERPOINT.docx")
    print("   - EDU_CONNECT_PRESENTATION.pptx (10 slides)")
    print("\n📥 Téléchargez-les via l'explorateur de fichiers d'Emergent !")
